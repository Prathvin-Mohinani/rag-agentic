import os
import json
import fitz
import torch
import pandas as pd
import xml.etree.ElementTree as ET
import logging
import time
 
from io import BytesIO
from concurrent.futures import ThreadPoolExecutor, as_completed
from dotenv import load_dotenv
from docx import Document as DocxDocument
from pptx import Presentation
 
from azure.storage.blob import BlobServiceClient
from azure.core.exceptions import ResourceNotFoundError
 
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFaceEmbeddings
 
 
# -------------------------------------------------
# CONFIG & LOGGING
# -------------------------------------------------
 
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s"
)
 
load_dotenv()
 
AZURE_CONNECTION_STRING = os.getenv("AZURE_STORAGE_CONNECTION_STRING")
CONTAINER_NAME = os.getenv("CONTAINER_NAME")
 
MAX_FILE_SIZE_MB = 20
MAX_WORKERS = 5
 
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
CATEGORY_PATH = os.path.join(BASE_DIR, "data", "category.json")
VECTOR_DB_PATH = os.path.join(BASE_DIR, "vector_store")
 
# ✅ NEW (Minimal Addition): Prefixes for incremental blob ingestion + move
UNPROCESSED_PREFIX = os.getenv("UNPROCESSED_PREFIX", "unprocessed/")
PROCESSED_PREFIX = os.getenv("PROCESSED_PREFIX", "ragdocument/")
 
 
# -------------------------------------------------
# LOAD CATEGORY
# -------------------------------------------------
 
with open(CATEGORY_PATH, "r") as f:
    CATEGORY_SCHEMA = json.load(f)
 
HINTS = CATEGORY_SCHEMA["hints"]
DOMAIN_DEFAULTS = CATEGORY_SCHEMA["domain_mapping_defaults"]
 
 
# -------------------------------------------------
# CLASSIFICATION
# -------------------------------------------------
 
def classify_document(text, filename):
    text_lower = text.lower()
    filename_lower = filename.lower()
 
    for label, hint in HINTS.items():
        filename_rules = [str(x).lower() for x in hint.get("filename_contains", [])]
        keyword_rules = [str(x).lower() for x in hint.get("keywords", [])]
        field_rules = [str(x).lower() for x in hint.get("field_markers", [])]
 
        if (
            any(rule in filename_lower for rule in filename_rules)
            or any(rule in text_lower for rule in keyword_rules)
            or any(rule in text_lower for rule in field_rules)
        ):
            group = label.split("/")[0]
            return {
                "label": label,
                "category": label.split("/")[-1],
                "domain": DOMAIN_DEFAULTS.get(group, "General"),
            }
 
    return {"label": "Unknown", "category": "Unknown", "domain": "General"}
 
 
# -------------------------------------------------
# FILE READERS
# -------------------------------------------------
 
def read_pdf(file_bytes):
    text = ""
    doc = fitz.open(stream=file_bytes, filetype="pdf")
 
    for page in doc:
        page_text = page.get_text()
 
        if not page_text.strip():
            try:
                textpage = page.get_textpage_ocr()
                page_text = page.get_text("text", textpage=textpage)
            except Exception as e:
                logging.warning(f"OCR failed: {e}")
 
        text += page_text
 
    return text
 
 
def read_excel(file_bytes):
    """
    Advanced Excel parser for messy real-world files
    """
    try:
        df = pd.read_excel(BytesIO(file_bytes), header=None)
        df = df.dropna(how="all")
 
        text_blocks = []
        current_headers = None
 
        for _, row in df.iterrows():
            values = [str(x).strip() for x in row if str(x).strip() != "nan"]
 
            if not values:
                continue
 
            # Detect header row
            if any("business" in v.lower() for v in values):
                current_headers = values
                continue
 
            if current_headers:
                row_data = []
 
                for i in range(min(len(current_headers), len(values))):
                    row_data.append(f"{current_headers[i]}: {values[i]}")
 
                if row_data:
                    text_blocks.append(" | ".join(row_data))
 
        # fallback if no header found
        if not text_blocks:
            df = pd.read_excel(BytesIO(file_bytes))
            df = df.dropna(how="all").fillna(method="ffill")
 
            df.columns = [str(col).strip() for col in df.columns]
 
            for _, row in df.iterrows():
                row_data = [
                    f"{col}: {str(row[col]).strip()}"
                    for col in df.columns
                    if str(row[col]).strip().lower() != "nan"
                ]
                if row_data:
                    text_blocks.append(" | ".join(row_data))
 
        return "\n".join(text_blocks)
 
    except Exception as e:
        logging.error(f"Excel parsing failed: {e}")
        return ""
 
 
def read_docx(file_bytes):
    doc = DocxDocument(BytesIO(file_bytes))
    return "\n".join([p.text for p in doc.paragraphs])
 
 
def read_ppt(file_bytes):
    prs = Presentation(BytesIO(file_bytes))
    text = []
 
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
 
    return "\n".join(text)
 
 
def read_xml(file_bytes):
    try:
        tree = ET.parse(BytesIO(file_bytes))
        root = tree.getroot()
 
        return "\n".join(
            elem.text.strip()
            for elem in root.iter()
            if elem.text and elem.text.strip()
        )
 
    except Exception:
        return file_bytes.decode("utf-8", errors="ignore")
 
 
# -------------------------------------------------
# PROCESS SINGLE FILE (EXISTING)
# -------------------------------------------------
 
def process_blob(blob, container_client):
    try:
        if blob.size > MAX_FILE_SIZE_MB * 1024 * 1024:
            logging.warning(f"Skipping large file: {blob.name}")
            return None
 
        blob_client = container_client.get_blob_client(blob)
        file_bytes = blob_client.download_blob().readall()
 
        file_name = blob.name
        ext = file_name.split(".")[-1].lower()
 
        if ext == "pdf":
            text = read_pdf(file_bytes)
        elif ext in ["xlsx", "xls"]:
            text = read_excel(file_bytes)
        elif ext == "docx":
            text = read_docx(file_bytes)
        elif ext == "pptx":
            text = read_ppt(file_bytes)
        elif ext == "xml":
            text = read_xml(file_bytes)
        else:
            text = file_bytes.decode("utf-8", errors="ignore")
 
        metadata = classify_document(text, file_name)
 
        return Document(
            page_content=text,
            metadata={
                "source": file_name,
                "category": metadata["category"],
                "domain": metadata["domain"],
                "label": metadata["label"],
            },
        )
 
    except Exception as e:
        logging.error(f"Failed processing {blob.name}: {e}")
        return None
 
 
# -------------------------------------------------
# LOAD DOCUMENTS (BULK – EXISTING)
# -------------------------------------------------
 
def load_documents():
    docs = []
 
    blob_service = BlobServiceClient.from_connection_string(AZURE_CONNECTION_STRING)
    container_client = blob_service.get_container_client(CONTAINER_NAME)
 
    blobs = list(container_client.list_blobs())
 
    with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
        futures = [
            executor.submit(process_blob, blob, container_client)
            for blob in blobs
        ]
 
        for future in as_completed(futures):
            result = future.result()
            if result:
                docs.append(result)
 
    logging.info(f"Total Documents Loaded: {len(docs)}")
    return docs
 
 
# -------------------------------------------------
# CHUNKING (EXISTING)
# -------------------------------------------------
 
def chunk_documents(documents):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=100
    )
 
    chunks = splitter.split_documents(documents)
    logging.info(f"Total Chunks Created: {len(chunks)}")
 
    return chunks
 
 
# -------------------------------------------------
# VECTOR DB (BULK – EXISTING)
# -------------------------------------------------
 
def create_vector_db(chunks):
    device = "cuda" if torch.cuda.is_available() else "cpu"
    logging.info(f"Using device: {device}")
 
    embeddings = HuggingFaceEmbeddings(
        model_name="BAAI/bge-base-en-v1.5",
        model_kwargs={"device": device},
        encode_kwargs={"batch_size": 64},
    )
 
    batch_size = 1000
    db = None
 
    for i in range(0, len(chunks), batch_size):
        batch = chunks[i: i + batch_size]
        logging.info(f"Processing batch {i // batch_size + 1}")
 
        if db is None:
            db = FAISS.from_documents(batch, embeddings)
        else:
            db.add_documents(batch)
 
    db.save_local(VECTOR_DB_PATH)
    logging.info("FAISS DB Created Successfully")
 
 
# -------------------------------------------------
# ✅ INCREMENTAL INGESTION (EXISTING - YOU ALREADY HAVE)
# -------------------------------------------------
 
def ingest_single_document(document: Document):
    """
    Incrementally add ONE document into the existing FAISS vector store
    without rebuilding it.
    """
    device = "cuda" if torch.cuda.is_available() else "cpu"
    logging.info(f"Ingesting single document using device: {device}")
 
    embeddings = HuggingFaceEmbeddings(
        model_name="BAAI/bge-base-en-v1.5",
        model_kwargs={"device": device},
        encode_kwargs={"batch_size": 64},
    )
 
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=100
    )
 
    chunks = splitter.split_documents([document])
 
    if os.path.exists(VECTOR_DB_PATH):
        db = FAISS.load_local(
            VECTOR_DB_PATH,
            embeddings,
            allow_dangerous_deserialization=True,
        )
        db.add_documents(chunks)
        logging.info("Existing FAISS DB updated")
    else:
        db = FAISS.from_documents(chunks, embeddings)
        logging.info("New FAISS DB created")
 
    db.save_local(VECTOR_DB_PATH)
    logging.info("FAISS DB saved successfully")
 
 
# -------------------------------------------------
# ✅ MINIMAL ADDITIONS (NEW): Move + Ingest by blob name
# -------------------------------------------------
 
def move_blob_within_container(container_client, source_name: str, dest_name: str, timeout_sec: int = 120):
    """
    Production-safe move: copy -> poll copy status -> delete source.
    """
    source_blob = container_client.get_blob_client(source_name)
    dest_blob = container_client.get_blob_client(dest_name)
 
    logging.info(f"Copy starting: {source_name} -> {dest_name}")
    dest_blob.start_copy_from_url(source_blob.url)
 
    # Poll copy status
    start = time.time()
    while True:
        props = dest_blob.get_blob_properties()
        copy_status = props.copy.status if props.copy else None
 
        if copy_status == "success":
            break
        if copy_status in ("failed", "aborted"):
            raise RuntimeError(f"Copy failed for {dest_name} with status={copy_status}")
 
        if time.time() - start > timeout_sec:
            raise TimeoutError(f"Copy timed out for {dest_name}")
 
        time.sleep(1)
 
    # Delete source after successful copy
    source_blob.delete_blob()
    logging.info(f"Moved blob: {source_name} -> {dest_name}")
 
 
def ingest_blob_by_name(blob_name: str):
    """
    Wrapper for Azure Functions QueueIndexer:
    - expects blob_name like 'unprocessed/2026-04-29/file.pdf'
    - reads blob properties
    - uses existing process_blob() to create Document
    - calls ingest_single_document()
    - moves blob to processed prefix after success
    """
    if not blob_name.startswith(UNPROCESSED_PREFIX):
        logging.warning(f"Skipping blob not under '{UNPROCESSED_PREFIX}': {blob_name}")
        return
 
    blob_service = BlobServiceClient.from_connection_string(AZURE_CONNECTION_STRING)
    container_client = blob_service.get_container_client(CONTAINER_NAME)
 
    try:
        blob_client = container_client.get_blob_client(blob_name)
        props = blob_client.get_blob_properties()  # has .name and .size
 
        document = process_blob(props, container_client)
        if not document:
            logging.warning(f"process_blob returned None for: {blob_name}")
            return
 
        # Incremental FAISS update
        ingest_single_document(document)
 
        # Move blob to processed prefix (preserve subfolders)
        relative = blob_name[len(UNPROCESSED_PREFIX):]  # everything after 'unprocessed/'
        dest_name = f"{PROCESSED_PREFIX}{relative}"
 
        move_blob_within_container(container_client, blob_name, dest_name)
 
    except ResourceNotFoundError:
        logging.error(f"Blob not found (maybe already moved): {blob_name}")
    except Exception as e:
        logging.error(f"Failed ingest_blob_by_name for {blob_name}: {e}")
        raise
 
 
# -------------------------------------------------
# MAIN (BULK MODE – OPTIONAL / EXISTING)
# -------------------------------------------------
 
if __name__ == "__main__":
    logging.info("Starting ingestion pipeline...")
    docs = load_documents()
    chunks = chunk_documents(docs)
    create_vector_db(chunks)
    logging.info("Ingestion completed successfully")